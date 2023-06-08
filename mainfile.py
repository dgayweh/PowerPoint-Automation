from pptx import Presentation
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.chart import  XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LABEL_POSITION
import os


# Load the data from the CSV file
df = pd.read_csv("sampledata.csv")
df = df.astype(str)
df=df.dropna()

title_text=''
# set the font and size of the body text
# set the font and size of the body text
font_name = ""
font_size = Pt(14)
logo_file=""
# create a new presentation
prs = Presentation()

# add a title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
#background color
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(52, 79, 158 )
#Title Text formatting
title = slide.shapes.title
title.text = title_text
title.text_frame.paragraphs[0].font.name = font_name
title.text_frame.paragraphs[0].font.size = Inches(0.555556)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


#Text Lines
line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(3.5), Inches(9.5), Inches(3.5))
line1.line.color.rgb = RGBColor(210, 149, 0 )
line1.line.width = Inches(0.02)
line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(3.55), Inches(9.5), Inches(3.55))
line2.line.color.rgb = RGBColor(210, 149, 0 )
line2.line.width = Inches(0.01)

# add the company logo to the slide
logo_path = logo_file
logo = slide.shapes.add_picture(logo_path, Inches(8), Inches(0.5), Inches(1.5), Inches(0.75))

#Next Slide 
content=""
slide = prs.slides.add_slide(prs.slide_layouts[1])

# add title to the slide
title = slide.shapes.title
title.text = "Introduction"
title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(52, 79, 158 )
title.text_frame.paragraphs[0].font.name = font_name

# add content to the slide

body = slide.placeholders[1]
body.left = Inches(0.5)
body.top = Inches(1.5)
body.width = Inches(9)
body.height = Inches(5.5)
body.text=content



for paragraph in body.text_frame.paragraphs:
    paragraph.bullet = None
    
for paragraph in body.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = font_size
        paragraph.bullet = None

logo_path = logo_file
logo = slide.shapes.add_picture(logo_path, Inches(8), Inches(0.5), Inches(1.5), Inches(0.75))
line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
line.line.color.rgb = RGBColor(4, 37, 58)
line.line.width = Inches(0.05)
line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(1.45), Inches(9.5), Inches(1.45))
line.line.color.rgb = RGBColor(210, 149, 0 )
line.line.width = Inches(0.02)
# save the presentation





for column in df.columns:
    # Get the unique values in the column
    unique_values = df[column].nunique(dropna=True)
    # Create a slide for the column
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Add a title to the slide with the column name

    if unique_values <=2:
    # Create the data for the pie chart
        data = df[column].dropna().value_counts().to_dict()  # exclude empty cells
        chart_data = ChartData()
        chart_data.categories = list(filter(lambda x: x != 'nan', data.keys()))
        chart_data.add_series(column, list(data.values()))
        # chart_data(number_format='0%')

        # Add the pie chart to the slide
        x, y, cx, cy = Inches(0.1), Inches(2), Inches(5), Inches(4)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart
        series = chart.series[0]
        data_labels = series.data_labels
        data_labels.show_value = True

        # Show data labels
        data_labels = series.data_labels
        data_labels.show_value = True
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        #Color
        #First Color is always blue
        chart.series[0].points[0].format.fill.solid()
        chart.series[0].points[0].format.fill.fore_color.rgb = RGBColor(52, 79, 158)
        # Set the fill color of the second slice to gold/grey
        chart.series[0].points[1].format.fill.solid()
        chart.series[0].points[1].format.fill.fore_color.rgb = RGBColor(166, 166, 166)
        # chart.series[0].points[1].format.fill.fore_color.rgb = RGBColor(210, 149, 0)
        logo_path = logo_file
        logo = slide.shapes.add_picture(logo_path, Inches(8), Inches(0.5), Inches(1.5), Inches(0.75))
        logo = slide.shapes.add_picture(logo_path, Inches(8), Inches(0.5), Inches(1.5), Inches(0.75))
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
        line.line.color.rgb = RGBColor(4, 37, 58)
        line.line.width = Inches(0.05)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(1.45), Inches(9.5), Inches(1.45))
        line.line.color.rgb = RGBColor(210, 149, 0 )
        line.line.width = Inches(0.02)
    # If there are more than two unique values, create a bar chart
    else:
        # Create the data for the bar chart
        data = df[column].dropna().value_counts().to_dict()
        chart_data = ChartData()
        chart_data.categories = list(filter(lambda x: x != 'nan', data.keys()))
        chart_data.add_series(column, list(data.values()))

        # Add the bar chart to the slide
        x, y, cx, cy = Inches(0.1), Inches(2), Inches(5), Inches(4)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            
        ).chart
        series = chart.series[0]
        data_labels = series.data_labels
        data_labels.show_value = True

        # Show data labels
        data_labels = series.data_labels
        data_labels.show_value = True
            
        chart.series[0].points[0].format.fill.solid()
        chart.series[0].points[0].format.fill.fore_color.rgb = RGBColor(52, 79, 158)
        # Remove the horizontal gridlines and keep only the major vertical gridline
       # Set the line format of the horizontal gridlines to "No line"
        chart.category_axis.major_gridlines.format.line.color.rgb = RGBColor(54, 69, 79)
        chart.category_axis.major_gridlines.format.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
        chart.category_axis.major_gridlines.format.line.width = 1
        # Set the line format of the vertical gridlines to "No line"
        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(255, 255, 255)
        chart.value_axis.major_gridlines.format.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
        chart.value_axis.major_gridlines.format.line.width = 0
       
        value_axis = chart.value_axis
        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.minor_tick_mark = XL_TICK_MARK.NONE
        value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
        logo_path = logo_file
        logo = slide.shapes.add_picture(logo_path, Inches(8), Inches(0.5), Inches(1.5), Inches(0.75))
        logo = slide.shapes.add_picture(logo_path, Inches(8), Inches(0.5), Inches(1.5), Inches(0.75))
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
        line.line.color.rgb = RGBColor(4, 37, 58)
        line.line.width = Inches(0.05)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.5), Inches(1.45), Inches(9.5), Inches(1.45))
        line.line.color.rgb = RGBColor(210, 149, 0 )
        line.line.width = Inches(0.02)
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = RGBColor(52, 79, 158 )
        chart.font.size = Pt(10)
        
prs.save("powerpoint.pptx")
os.startfile("power.pptx")
